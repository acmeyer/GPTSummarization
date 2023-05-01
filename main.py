import argparse
import openai
import os
import tiktoken
import requests
import pickle
import mimetypes
from tenacity import retry, wait_random_exponential, stop_after_attempt
from langchain import PromptTemplate
from langchain.embeddings import OpenAIEmbeddings
from langchain.text_splitter import TokenTextSplitter
from langchain.vectorstores.faiss import FAISS
from langchain.document_loaders import PyPDFLoader
from langchain.document_loaders.csv_loader import CSVLoader
from langchain.document_loaders import Docx2txtLoader
import pptx
from bs4 import BeautifulSoup
import dotenv
dotenv.load_dotenv('.env')

openai.api_key = os.getenv("OPENAI_API_KEY")
GPT_4_MODEL = "gpt-4"
GPT_3_MODEL = "gpt-3.5-turbo"
DEFAULT_CHAT_MODEL = GPT_3_MODEL
MODEL_TEMPERATURE = 0.0
OPENAI_EMBEDDING_ENCODING = "cl100k_base" # this the encoding for text-embedding-ada-002
MAX_TOKENS = 500
CHUNK_TOKEN_SIZE = 200

# check if data directory exists
if not os.path.exists("data/summaries"):
    os.makedirs("data/summaries")

def count_tokens(text: str) -> int:
    """Returns the number of tokens in the given text."""
    encoding = tiktoken.get_encoding(OPENAI_EMBEDDING_ENCODING)
    tokens = encoding.encode(text)
    num_tokens = len(tokens)
    return num_tokens, tokens

@retry(wait=wait_random_exponential(min=1, max=20), stop=stop_after_attempt(6))
def get_chat_completion(prompt: str, messages: list[dict], model: str = 'gpt-3.5-turbo'):
    """Returns ChatGPT's response to the given prompt."""
    system_message = [{"role": "system", "content": prompt}]
    conversation_messages = system_message + messages
    response = openai.ChatCompletion.create(
        model=model,
        messages=conversation_messages,
        temperature=MODEL_TEMPERATURE
    )
    return response.choices[0]['message']['content'].strip()

def get_summary(full_text: str, summary_type: str, model: str = 'gpt-3.5-turbo') -> str:
    with open("prompts/summarization_prompt.md") as f:
        template = f.read()
    prompt = PromptTemplate(
        input_variables=[],
        template=template,
    ).format()
    conversation_messages = []
    if summary_type == "detailed":
        user_input = f"Here's the content you should summarize:\n\n{full_text}\n\n----\n\nI would like you to produce a detailed summary of this content."
    else:
        user_input = f"Here's the content you should summarize:\n\n{full_text}\n\n----\n\nI would like you to produce a short summary of this content. It should be a few sentences at most."
    conversation_messages.append({"role": "user", "content": user_input})
    return get_chat_completion(prompt, conversation_messages, model=model)

def extract_text_from_url(url: str) -> str:
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.36'
        }
        response = requests.get(url, headers=headers)
        response.raise_for_status()  # Raise an exception if the request returned an HTTP error
        soup = BeautifulSoup(response.text, 'html.parser')
        text = soup.get_text()
        return text
    except requests.exceptions.RequestException as e:
        print(f"An error occurred while fetching the URL: {e}")
        return None
    
def write_to_file(text: str, filename: str):
    with open(filename, 'w') as f:
        f.write(text)
    print(f"Summary saved to {filename}")


def extract_text_from_filepath(filepath: str) -> str:
    mimetype, _ = mimetypes.guess_type(filepath)
    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise Exception("Unsupported file type")
    if filepath.endswith(".md"):
        mimetype = "text/markdown"
    if mimetype == "application/pdf":
        # Extract text from pdf
        loader = PyPDFLoader(filepath)
        pages = loader.load_and_split()
        extracted_text = " ".join([page.page_content for page in pages])
    elif mimetype == "text/plain" or mimetype == "text/markdown":
        # Read text from plain text file
        with open(filepath, "rb") as file:
            extracted_text = file.read().decode("utf-8")
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # Extract text from docx
        loader = Docx2txtLoader(filepath)
        data = loader.load()
        extracted_text = " ".join([page.page_content for page in data])
    elif mimetype == "text/csv":
        # Extract text from csv
        loader = CSVLoader(file_path=filepath)
        data = loader.load()
        for row in data:
            extracted_text += " ".join(row.page_content) + "\n"
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # Extract text from pptx
        extracted_text = ""
        with open(filepath, "rb") as file:
            presentation = pptx.Presentation(file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                extracted_text += run.text + " "
                        extracted_text += "\n"
    else:
        # Unsupported file type
        raise ValueError("Unsupported file type: {}".format(mimetype))

    return extracted_text



def create_summary(url: str, filepath: str, summary_type: str, model: str = DEFAULT_CHAT_MODEL) -> str:
    if filepath:
        full_text = extract_text_from_filepath(filepath)
        num_tokens, _ = count_tokens(full_text)
        title = filepath.split("/")[-1]
        title = title.replace(" ", "_")
    elif url:
        full_text = extract_text_from_url(url)
        num_tokens, _ = count_tokens(full_text)
        title = url.split("/")[-1]
        if not title or title == '' or title == 'index.html':
            title = url.split("/")[-2]

    embeddings = OpenAIEmbeddings()
    texts_for_embeddings = []
    if num_tokens > MAX_TOKENS:
        text_splitter = TokenTextSplitter(chunk_size=CHUNK_TOKEN_SIZE, chunk_overlap=20)
        texts = text_splitter.split_text(full_text)

        summaries = []
        for text in texts:
            summary = get_summary(text, summary_type="short", model=model)
            summaries.append(summary)
            texts_for_embeddings.append(text)
        all_summaries = "\n\n".join(summaries)
        full_response = get_summary(all_summaries, summary_type="detailed", model=GPT_4_MODEL)
        write_to_file(full_response, f"data/summaries/{title}.txt")
        response = full_response
    else:
        response = get_summary(full_text, summary_type="detailed", model=GPT_4_MODEL)
        write_to_file(response, f"data/summaries/{title}.txt")

    print(f"Saving vectorstore into data/{title}_vectorstore.pkl")
    vectorstore = FAISS.from_texts(texts_for_embeddings, embeddings)
    with open(f"data/{title}_vectorstore.pkl", "wb") as f:
        pickle.dump(vectorstore, f)

    return response, title


# To run directly on command line, run `python main.py`
if __name__ == "__main__":
    argparser = argparse.ArgumentParser()
    argparser.add_argument("-u", "--url", type=str, help="URL to summarize")
    argparser.add_argument("-f", "--filepath", type=str, help="Path to file to summarize")
    argparser.add_argument("-m", "--model", type=str, help="Model to use", default="gpt-3.5-turbo")
    argparser.add_argument("-s", "--summary_type", type=str, help="Type of summary to generate", default="short")
    args = argparser.parse_args()
    url = args.url
    summary_type = args.summary_type
    filepath = args.filepath
    model = args.model

    assert url or filepath, "Must provide either a URL or a filepath"

    response, title = create_summary(url, filepath, summary_type, model)
    print(f'\033[96m\033[1mSummary: {response}\033[0m\033[1m')

    conversation_messages = []
    with open("prompts/chat_prompt.md", "r") as f:
        template = f.read()
    prompt = PromptTemplate(
        input_variables=[],
        template=template,
    ).format()
    with open(f"data/{title}_vectorstore.pkl", "rb") as f:
        vectorstore = pickle.load(f)
    while (user_input := input('You: ').strip()) != "":
        relevant_texts = vectorstore.similarity_search(user_input, top_k=3)
        relevant_content = "\n\n".join([f"{text}" for text in relevant_texts])
        user_content = f"""
        Relevant content:
        {relevant_content}
        =====
        {user_input}
        """
        user_message = {"role": "user", "content": user_content}
        conversation_messages.append(user_message)
        answer = get_chat_completion(prompt, messages=conversation_messages, model=DEFAULT_CHAT_MODEL)
        conversation_messages.append({"role": "assistant", "content": answer})
        print(f'\033[96m\033[1mGPT: {answer}\033[0m\033[1m')
